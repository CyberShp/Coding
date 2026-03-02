#!/usr/bin/env python3
"""Generate the Pangu Gray/White-box Testing Project Plan PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Color palette ---
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xA0)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
MID_TEXT = RGBColor(0x4B, 0x55, 0x63)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
HEADER_BG = RGBColor(0xEF, 0xF6, 0xFF)
ROW_ALT = RGBColor(0xF9, 0xFB, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_para(tf, text, font_size=14, color=DARK_TEXT, bold=False, space_before=0, space_after=0,
             alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def set_cell_text(cell, text, font_size=11, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_fill(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    return table


def add_slide_number(slide, num, total=16):
    add_text_box(slide, SW - Inches(1.5), SH - Inches(0.5), Inches(1.3), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=MID_TEXT, alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, section_text, slide_num=None, total=16):
    add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), fill_color=ACCENT_BLUE)
    if slide_num:
        add_slide_number(slide, slide_num, total)


# ============================================================
# SLIDE 1 — Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape(slide, Inches(0), Inches(0), SW, Inches(0.08), fill_color=ACCENT_TEAL)
add_shape(slide, Inches(0), SH - Inches(0.08), SW, Inches(0.08), fill_color=ACCENT_TEAL)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "盘古平台灰白盒测试系统建设项目", font_size=40, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.8),
             "精准测试 × AI 赋能 — 构建代码感知的测试专业能力", font_size=22, color=ACCENT_TEAL)

line = add_shape(slide, Inches(1.5), Inches(4.2), Inches(3), Inches(0.04), fill_color=ACCENT_TEAL)

add_text_box(slide, Inches(1.5), Inches(4.8), Inches(6), Inches(0.5),
             "项目周期：2026年 Q2 — Q3（6个月）", font_size=16, color=RGBColor(0xA0, 0xAE, 0xC0))
add_text_box(slide, Inches(1.5), Inches(5.3), Inches(6), Inches(0.5),
             "试点特性：NVMe over TCP  |  KV", font_size=16, color=RGBColor(0xA0, 0xAE, 0xC0))
add_text_box(slide, Inches(1.5), Inches(5.8), Inches(6), Inches(0.5),
             "日期：2026年3月", font_size=16, color=RGBColor(0xA0, 0xAE, 0xC0))

# ============================================================
# SLIDE 2 — Background & Problems
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "背景", 2)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "背景与问题分析", font_size=28, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
             "从黑盒盲区到灰白盒精准打击", font_size=14, color=MID_TEXT)

problems = [
    ("代码异常分支漏测", "错误处理路径、资源释放路径未被覆盖\n逃逸到客户现场", ACCENT_RED),
    ("流程并发问题", "多线程竞态、锁序异常\n异步时序偶发性故障", ACCENT_ORANGE),
    ("协议报文漏测", "NVMe-oF TCP / KV 边界值\n异常报文场景覆盖不足", ACCENT_BLUE),
    ("测试设计依赖经验", "用例设计质量因人而异\n缺乏系统化方法论", MID_TEXT),
]

for i, (title, desc, color) in enumerate(problems):
    x = Inches(0.8 + i * 3.05)
    card = add_rounded_rect(slide, x, Inches(1.5), Inches(2.8), Inches(2.5),
                            fill_color=LIGHT_BG, line_color=BORDER_GRAY)
    top_bar = add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(0.06), fill_color=color)
    add_text_box(slide, x + Inches(0.2), Inches(1.75), Inches(2.4), Inches(0.5),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.5),
                 desc, font_size=12, color=MID_TEXT)

quote_box = add_rounded_rect(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.1),
                              fill_color=RGBColor(0xFF, 0xF7, 0xED), line_color=ACCENT_ORANGE)
add_text_box(slide, Inches(1.1), Inches(4.45), Inches(11), Inches(0.9),
             "核心矛盾：黑盒测试不知道代码里哪些函数、哪些故障处理分支会在同一时刻交汇，\n"
             "只能靠 N 次测试提高命中率 —— 灰白盒测试的核心就是精准找到多函数交汇临界点。",
             font_size=14, color=RGBColor(0x92, 0x40, 0x0E), bold=False)

root_box = add_rounded_rect(slide, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.4),
                             fill_color=RGBColor(0xEF, 0xF6, 0xFF), line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(5.75), Inches(11), Inches(0.4),
             "根因分析", font_size=14, color=ACCENT_BLUE, bold=True)
causes = [
    "黑盒测试无法感知代码内部分支结构 → 异常分支漏测",
    "缺乏对共享变量和并发路径的系统性分析 → 并发问题逃逸",
    "协议实现与代码流程的关联分析缺位 → 协议报文覆盖不足",
]
for j, c in enumerate(causes):
    add_text_box(slide, Inches(1.3), Inches(6.1 + j * 0.3), Inches(10.5), Inches(0.3),
                 f"▸ {c}", font_size=11, color=MID_TEXT)

# ============================================================
# SLIDE 3 — Three Goals
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "目标", 3)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "项目三大目标", font_size=28, color=DARK_TEXT, bold=True)

goals = [
    ("01", "流程建设", "建立盘古平台灰白盒测试系统流程",
     "版本交付各阶段有明确的灰白盒测试活动\n出入口标准、关键产出", ACCENT_BLUE),
    ("02", "质量改进", "漏测问题数改进 ≥ 30%",
     "代码异常分支、流程并发、协议报文\n三类漏测以历史数据为基线量化改进", ACCENT_TEAL),
    ("03", "能力突破", "精准测试与 AI 模型结合",
     "代码调用链辅助分析、业务流程生成\n异步时序偶发性问题拦截", ACCENT_ORANGE),
]

for i, (num, title, subtitle, desc, color) in enumerate(goals):
    x = Inches(0.8 + i * 4.0)
    card = add_rounded_rect(slide, x, Inches(1.4), Inches(3.7), Inches(5.2),
                            fill_color=LIGHT_BG, line_color=color)
    num_box = add_rounded_rect(slide, x + Inches(0.2), Inches(1.65), Inches(0.8), Inches(0.8),
                                fill_color=color)
    add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(0.8), Inches(0.8),
                 num, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(1.2), Inches(1.75), Inches(2.3), Inches(0.5),
                 title, font_size=22, color=color, bold=True)
    add_text_box(slide, x + Inches(1.2), Inches(2.25), Inches(2.3), Inches(0.5),
                 subtitle, font_size=13, color=DARK_TEXT, bold=True)
    add_shape(slide, x + Inches(0.3), Inches(2.95), Inches(3.1), Inches(0.02), fill_color=BORDER_GRAY)
    add_text_box(slide, x + Inches(0.3), Inches(3.2), Inches(3.1), Inches(2.5),
                 desc, font_size=12, color=MID_TEXT)

arrow_y = Inches(6.85)
add_shape(slide, Inches(0.8), arrow_y, Inches(11.7), Inches(0.04), fill_color=ACCENT_BLUE)
labels = ["流程是基座", "质量是核心", "AI是拔高"]
for i, lb in enumerate(labels):
    add_text_box(slide, Inches(1.8 + i * 4.0), arrow_y + Inches(0.1), Inches(2), Inches(0.4),
                 f"▲ {lb}", font_size=11, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4 — Metrics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "度量", 4)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "目标衡量体系（定量化）", font_size=28, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
             "每项目标配套可量化指标，用数据说话", font_size=14, color=MID_TEXT)

headers = ["目标", "衡量指标", "基线", "目标值", "数据来源", "频率"]
col_ws = [Inches(1.1), Inches(3.0), Inches(1.8), Inches(1.8), Inches(2.2), Inches(1.5)]
data = [
    ["流程建设", "灰白盒流程文档发布", "0", "≥1 套完整流程", "项目交付物", "里程碑"],
    ["", "试点特性活动执行率", "0%", "≥90%", "版本测试报告", "每版本"],
    ["", "DT用例中灰白盒占比", "统计当前", "提升至≥15%", "用例管理系统", "每版本"],
    ["质量改进", "异常分支类漏测问题数", "历史均值", "下降 ≥30%", "缺陷管理系统", "每版本"],
    ["", "并发类漏测问题数", "历史均值", "下降 ≥30%", "缺陷管理系统", "每版本"],
    ["", "协议报文类漏测问题数", "历史均值", "下降 ≥30%", "缺陷管理系统", "每版本"],
    ["能力突破", "AI辅助分析特性数", "0", "≥2 个", "分析报告", "项目结束"],
    ["", "调用链→测试场景转化用例", "0", "≥20 条", "用例管理系统", "项目结束"],
    ["", "覆盖率AI分析报告", "0", "≥2 份", "分析报告", "项目结束"],
]

table = add_table(slide, len(data) + 1, len(headers),
                  Inches(0.6), Inches(1.35), Inches(11.4), Inches(5.5))

for ci, h in enumerate(headers):
    table.columns[ci].width = col_ws[ci]
    set_cell_text(table.cell(0, ci), h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    set_cell_fill(table.cell(0, ci), ACCENT_BLUE)

goal_colors = {
    "流程建设": RGBColor(0xE0, 0xF2, 0xFE),
    "质量改进": RGBColor(0xD1, 0xFA, 0xE5),
    "能力突破": RGBColor(0xFF, 0xED, 0xD5),
}
current_goal_color = None
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        set_cell_text(table.cell(ri + 1, ci), val, font_size=10, alignment=PP_ALIGN.CENTER)
        if row[0] in goal_colors:
            current_goal_color = goal_colors[row[0]]
        if ci == 0 and row[0]:
            set_cell_fill(table.cell(ri + 1, ci), current_goal_color)
            set_cell_text(table.cell(ri + 1, ci), val, font_size=10, bold=True, alignment=PP_ALIGN.CENTER)
        elif ri % 2 == 1:
            set_cell_fill(table.cell(ri + 1, ci), ROW_ALT)

note_box = add_rounded_rect(slide, Inches(0.6), Inches(6.95), Inches(11.4), Inches(0.4),
                             fill_color=RGBColor(0xFE, 0xF3, 0xC7))
add_text_box(slide, Inches(0.9), Inches(6.97), Inches(10.8), Inches(0.35),
             "💡 基线值需在M0阶段从缺陷管理系统导出确认 | 30%改进覆盖三个维度独立统计 | 每版本度量持续跟踪",
             font_size=10, color=RGBColor(0x92, 0x40, 0x0E))

# ============================================================
# SLIDE 5 — Solution Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "方案", 5)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "整体解决方案全景", font_size=28, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
             "三层架构：流程基座 → 特性覆盖 → AI赋能", font_size=14, color=MID_TEXT)

layers = [
    ("方案三：精准测试 × AI 赋能",
     "精准测试工具 + AI模型 → 覆盖率分析 → 调用链→场景转换 → 时序问题拦截\nGrayScope 辅助分析（持续验证迭代中）",
     ACCENT_ORANGE, RGBColor(0xFF, 0xF7, 0xED)),
    ("方案二：特性测试策略",
     "新增特性：抓包 + 设计文档 + 协议标准 → 协议实现与代码流程 → DT用例\n继承特性：SFMEA + 协议规范 → 协议分析仪 + 网损仪 → 边界值 + 异常分支",
     ACCENT_TEAL, RGBColor(0xEC, 0xFD, 0xF5)),
    ("方案一：流程与标准",
     "定义版本交付各阶段灰白盒测试活动 → 出入口标准 → 关键产出",
     ACCENT_BLUE, RGBColor(0xEF, 0xF6, 0xFF)),
]

for i, (title, desc, color, bg_color) in enumerate(layers):
    y = Inches(1.5 + i * 1.85)
    card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.65),
                            fill_color=bg_color, line_color=color)
    add_shape(slide, Inches(0.8), y, Inches(0.12), Inches(1.65), fill_color=color)
    add_text_box(slide, Inches(1.2), y + Inches(0.15), Inches(10), Inches(0.45),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.65), Inches(10.5), Inches(0.9),
                 desc, font_size=13, color=MID_TEXT)

for i in range(2):
    ay = Inches(3.15 + i * 1.85)
    arrow_shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                          Inches(6.3), ay, Inches(0.7), Inches(0.45))
    arrow_shape.fill.solid()
    arrow_shape.fill.fore_color.rgb = BORDER_GRAY
    arrow_shape.line.fill.background()
    arrow_shape.rotation = 180.0

rel_box = add_rounded_rect(slide, Inches(0.8), Inches(7.1 - 0.8), Inches(11.7), Inches(0.6),
                            fill_color=LIGHT_BG)
add_text_box(slide, Inches(1.1), Inches(7.1 - 0.75), Inches(11), Inches(0.5),
             "层级关系：方案一是基座（流程保障）→ 方案二是主体（特性覆盖）→ 方案三是拔高（工具与AI赋能）",
             font_size=12, color=MID_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6 — Process Framework
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "方案一", 6)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "方案一：灰白盒测试流程体系", font_size=28, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
             "嵌入版本交付流程，不增加独立阶段", font_size=14, color=MID_TEXT)

headers6 = ["版本阶段", "灰白盒活动", "入口标准", "出口标准", "关键产出"]
col6 = [Inches(1.4), Inches(2.8), Inches(2.2), Inches(2.6), Inches(2.4)]
data6 = [
    ["需求分析", "识别关键协议流程\n与代码模块", "特性需求基线化", "关键流程清单确认", "关键流程识别报告"],
    ["方案设计", "代码架构与协议\n实现走读", "设计文档基线化", "关键代码路径标记完成", "代码走读记录\n风险点清单"],
    ["编码阶段", "覆盖率采集配置\n增量分析", "代码提交", "增量代码覆盖率≥基线", "覆盖率报告"],
    ["DT 测试", "灰白盒用例设计\n与执行", "功能测试用例基线化", "灰白盒用例执行完成", "灰白盒测试报告"],
    ["SDV / SIT", "异常分支 & 并发\n场景专项测试", "DT测试出口达标", "异常/并发场景\n覆盖率达标", "专项测试报告"],
    ["版本发布", "灰白盒测试\n度量汇总", "所有阶段出口达标", "漏测改进数据达标", "版本灰白盒度量报告"],
]

table6 = add_table(slide, len(data6) + 1, len(headers6),
                   Inches(0.5), Inches(1.35), Inches(11.4), Inches(5.5))

for ci, h in enumerate(headers6):
    table6.columns[ci].width = col6[ci]
    set_cell_text(table6.cell(0, ci), h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    set_cell_fill(table6.cell(0, ci), ACCENT_BLUE)

stage_colors = [
    RGBColor(0xDB, 0xEA, 0xFE), RGBColor(0xD1, 0xFA, 0xE5), RGBColor(0xFE, 0xF9, 0xC3),
    RGBColor(0xFF, 0xED, 0xD5), RGBColor(0xFC, 0xE7, 0xF3), RGBColor(0xE0, 0xE7, 0xFF),
]
for ri, row in enumerate(data6):
    for ci, val in enumerate(row):
        set_cell_text(table6.cell(ri + 1, ci), val, font_size=10, alignment=PP_ALIGN.CENTER)
        if ci == 0:
            set_cell_fill(table6.cell(ri + 1, ci), stage_colors[ri])
            set_cell_text(table6.cell(ri + 1, ci), val, font_size=10, bold=True, alignment=PP_ALIGN.CENTER)
        elif ri % 2 == 1:
            set_cell_fill(table6.cell(ri + 1, ci), ROW_ALT)

# ============================================================
# SLIDE 7 — New Feature Strategy (NVMe-oF TCP)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "方案二-A", 7)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "方案二-A：新增特性测试策略", font_size=28, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
             "适用对象：NVMe over TCP 新增特性", font_size=14, color=ACCENT_TEAL, bold=True)

flow_steps = [
    ("设计文档", ACCENT_BLUE),
    ("抓包工具", ACCENT_TEAL),
    ("协议标准", ACCENT_ORANGE),
]
for i, (label, color) in enumerate(flow_steps):
    x = Inches(0.8)
    y = Inches(1.5 + i * 0.75)
    box = add_rounded_rect(slide, x, y, Inches(1.8), Inches(0.55), fill_color=color)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.08), Inches(1.6), Inches(0.4),
                 label, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.8), Inches(1.95),
                                 Inches(0.8), Inches(0.45))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = BORDER_GRAY
arrow1.line.fill.background()

mid_box = add_rounded_rect(slide, Inches(3.8), Inches(1.45), Inches(3.8), Inches(2.0),
                            fill_color=RGBColor(0xEF, 0xF6, 0xFF), line_color=ACCENT_BLUE)
add_text_box(slide, Inches(4.0), Inches(1.55), Inches(3.4), Inches(0.4),
             "协议实现与关键代码流程分析", font_size=13, color=ACCENT_BLUE, bold=True)
mid_items = ["关键函数调用链梳理", "错误处理分支标记", "并发路径识别", "多函数交汇临界点分析"]
for j, item in enumerate(mid_items):
    add_text_box(slide, Inches(4.0), Inches(2.0 + j * 0.33), Inches(3.4), Inches(0.3),
                 f"▸ {item}", font_size=11, color=MID_TEXT)

arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.8), Inches(1.95),
                                 Inches(0.8), Inches(0.45))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = BORDER_GRAY
arrow2.line.fill.background()

right_box = add_rounded_rect(slide, Inches(8.8), Inches(1.55), Inches(3.7), Inches(1.8),
                              fill_color=ACCENT_TEAL)
add_text_box(slide, Inches(9.0), Inches(1.7), Inches(3.3), Inches(0.4),
             "DT用例设计与执行", font_size=15, color=WHITE, bold=True)
add_text_box(slide, Inches(9.0), Inches(2.15), Inches(3.3), Inches(1.0),
             "灰盒用例围绕\n「多函数交汇临界点」\n一条用例验证一个交汇场景",
             font_size=12, color=RGBColor(0xD1, 0xFA, 0xE5))

headers7 = ["步骤", "活动", "工具/方法", "产出"]
col7 = [Inches(0.8), Inches(3.0), Inches(3.8), Inches(3.8)]
data7 = [
    ["1", "协议抓包与报文解析", "Wireshark / tcpdump / 协议分析仪", "协议交互流程图"],
    ["2", "设计文档走读", "设计文档 + 代码对照", "关键代码函数清单"],
    ["3", "代码流程分析", "精准测试工具（覆盖率）+ 代码走读", "调用链图、分支路径图"],
    ["4", "灰盒用例设计", "基于调用链 + 多函数交汇点", "灰盒测试用例集"],
    ["5", "参与DT用例评审", "灰盒视角补充用例", "DT用例中灰盒用例"],
    ["6", "执行与度量", "代码覆盖率对比", "覆盖率提升数据"],
]

table7 = add_table(slide, len(data7) + 1, len(headers7),
                   Inches(0.5), Inches(3.9), Inches(11.4), Inches(3.2))
for ci, h in enumerate(headers7):
    table7.columns[ci].width = col7[ci]
    set_cell_text(table7.cell(0, ci), h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    set_cell_fill(table7.cell(0, ci), ACCENT_TEAL)

for ri, row in enumerate(data7):
    for ci, val in enumerate(row):
        set_cell_text(table7.cell(ri + 1, ci), val, font_size=10, alignment=PP_ALIGN.CENTER)
        if ri % 2 == 1:
            set_cell_fill(table7.cell(ri + 1, ci), ROW_ALT)

# ============================================================
# SLIDE 8 — Inherited Feature Strategy (KV)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "方案二-B", 8)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "方案二-B：继承特性测试策略", font_size=28, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
             "适用对象：KV 继承特性", font_size=14, color=ACCENT_ORANGE, bold=True)

left_sources = [("SFMEA", ACCENT_RED), ("协议规范", ACCENT_BLUE)]
for i, (label, color) in enumerate(left_sources):
    y = Inches(1.6 + i * 0.75)
    box = add_rounded_rect(slide, Inches(0.8), y, Inches(1.8), Inches(0.55), fill_color=color)
    add_text_box(slide, Inches(0.9), y + Inches(0.08), Inches(1.6), Inches(0.4),
                 label, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

arrow_s8 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.8), Inches(1.95),
                                   Inches(0.8), Inches(0.45))
arrow_s8.fill.solid()
arrow_s8.fill.fore_color.rgb = BORDER_GRAY
arrow_s8.line.fill.background()

mid_box8 = add_rounded_rect(slide, Inches(3.8), Inches(1.45), Inches(3.8), Inches(2.0),
                             fill_color=RGBColor(0xFF, 0xF7, 0xED), line_color=ACCENT_ORANGE)
add_text_box(slide, Inches(4.0), Inches(1.55), Inches(3.4), Inches(0.4),
             "有效测试点提取", font_size=13, color=ACCENT_ORANGE, bold=True)
mid_items8 = ["历史缺陷模式匹配", "关键边界条件识别", "异常注入场景设计", "风险评分排序"]
for j, item in enumerate(mid_items8):
    add_text_box(slide, Inches(4.0), Inches(2.0 + j * 0.33), Inches(3.4), Inches(0.3),
                 f"▸ {item}", font_size=11, color=MID_TEXT)

arrow_s8b = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.8), Inches(1.95),
                                    Inches(0.8), Inches(0.45))
arrow_s8b.fill.solid()
arrow_s8b.fill.fore_color.rgb = BORDER_GRAY
arrow_s8b.line.fill.background()

right_items = [("协议分析仪", "报文级精确验证"), ("网损仪", "网络异常模拟")]
for i, (t, d) in enumerate(right_items):
    y = Inches(1.5 + i * 1.1)
    box = add_rounded_rect(slide, Inches(8.8), y, Inches(3.7), Inches(0.9),
                            fill_color=ACCENT_ORANGE)
    add_text_box(slide, Inches(9.0), y + Inches(0.05), Inches(3.3), Inches(0.4),
                 t, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.0), y + Inches(0.45), Inches(3.3), Inches(0.4),
                 d, font_size=11, color=RGBColor(0xFF, 0xED, 0xD5), alignment=PP_ALIGN.CENTER)

headers8 = ["步骤", "活动", "工具/方法", "产出"]
col8 = [Inches(0.8), Inches(3.0), Inches(3.8), Inches(3.8)]
data8 = [
    ["1", "SFMEA 分析", "SFMEA表 + 历史缺陷库", "失效模式与测试点映射"],
    ["2", "协议规范梳理", "协议标准文档", "边界值 & 异常场景清单"],
    ["3", "测试点优先级排序", "风险评分（影响×概率）", "分级测试点清单"],
    ["4", "测试设计与执行", "协议分析仪 + 网损仪", "异常报文 & 网络异常用例"],
    ["5", "覆盖率验证", "精准测试工具", "代码覆盖率报告"],
    ["6", "缺陷闭环", "缺陷管理系统", "新发现纳入知识库"],
]

table8 = add_table(slide, len(data8) + 1, len(headers8),
                   Inches(0.5), Inches(3.9), Inches(11.4), Inches(3.2))
for ci, h in enumerate(headers8):
    table8.columns[ci].width = col8[ci]
    set_cell_text(table8.cell(0, ci), h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    set_cell_fill(table8.cell(0, ci), ACCENT_ORANGE)

for ri, row in enumerate(data8):
    for ci, val in enumerate(row):
        set_cell_text(table8.cell(ri + 1, ci), val, font_size=10, alignment=PP_ALIGN.CENTER)
        if ri % 2 == 1:
            set_cell_fill(table8.cell(ri + 1, ci), ROW_ALT)

# ============================================================
# SLIDE 9 — Precise Testing x AI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "方案三", 9)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "方案三：精准测试 × AI 赋能", font_size=28, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
             "三项核心能力建设", font_size=14, color=MID_TEXT)

caps = [
    ("覆盖率 AI 分析", "代码覆盖率数据 + 代码结构",
     "识别高风险低覆盖区域\n关联历史缺陷模式",
     "风险热力图\n补充测试建议", ACCENT_BLUE),
    ("调用链 → 业务场景", "代码调用链 + 函数签名",
     "调用链语义理解\n映射到业务操作序列",
     "业务场景描述\n建议测试步骤", ACCENT_TEAL),
    ("异步时序分析", "并发代码结构 + 时序关系",
     "竞态条件识别\n时序窗口分析",
     "触发条件\n复现建议", ACCENT_ORANGE),
]

for i, (title, input_t, process_t, output_t, color) in enumerate(caps):
    x = Inches(0.8 + i * 4.05)
    card = add_rounded_rect(slide, x, Inches(1.35), Inches(3.8), Inches(3.6),
                            fill_color=LIGHT_BG, line_color=color)
    add_shape(slide, x, Inches(1.35), Inches(3.8), Inches(0.06), fill_color=color)
    add_text_box(slide, x + Inches(0.2), Inches(1.55), Inches(3.4), Inches(0.4),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    labels9 = [("输入", input_t), ("AI 处理", process_t), ("输出", output_t)]
    for j, (lb, val) in enumerate(labels9):
        yy = Inches(2.1 + j * 0.95)
        add_text_box(slide, x + Inches(0.2), yy, Inches(1.0), Inches(0.25),
                     lb, font_size=9, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), yy + Inches(0.22), Inches(3.4), Inches(0.65),
                     val, font_size=10, color=MID_TEXT)

gs_box = add_rounded_rect(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(2.0),
                           fill_color=RGBColor(0xFF, 0xF7, 0xED), line_color=ACCENT_ORANGE)
add_text_box(slide, Inches(1.1), Inches(5.3), Inches(6), Inches(0.4),
             "GrayScope 工具定位", font_size=15, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(1.1), Inches(5.7), Inches(4.5), Inches(0.3),
             "⚠ 辅助验证工具，持续迭代中", font_size=12, color=ACCENT_RED, bold=True)

gs_items = [
    ("当前状态", "九大分析模块已完成开发（分支路径/边界值/错误路径/调用图/并发风险等）"),
    ("项目角色", "辅助验证 — 对比验证人工分析结果的完整性和准确性"),
    ("纳入条件", "分析准确率 ≥70% 且减少人工分析时间 ≥30%"),
    ("风险预案", "不纳入也不影响核心目标，方案一二独立可运作"),
]
for j, (k, v) in enumerate(gs_items):
    add_text_box(slide, Inches(1.1), Inches(6.05 + j * 0.27), Inches(1.5), Inches(0.25),
                 k, font_size=10, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(2.6), Inches(6.05 + j * 0.27), Inches(9.5), Inches(0.25),
                 v, font_size=10, color=MID_TEXT)

# ============================================================
# SLIDE 10 — Pilot Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "试点", 10)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "试点特性技术方案", font_size=28, color=DARK_TEXT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(5.5), Inches(0.5),
             "NVMe over TCP — 关键流程", font_size=18, color=ACCENT_BLUE, bold=True)

h10a = ["关键流程", "灰白盒关注点", "测试方法"]
col10a = [Inches(1.8), Inches(2.5), Inches(1.9)]
d10a = [
    ["连接建立 Connect", "状态机异常分支\n超时处理、重连", "调用链分析 + 网损仪"],
    ["命令下发 I/O", "命令队列并发\n队列满异常处理", "并发路径 + 压力测试"],
    ["错误恢复", "错误处理函数交汇\n资源释放路径", "错误路径 + 故障注入"],
    ["断连处理", "并发I/O处理\n状态清理", "交汇分析 + 时序测试"],
]
t10a = add_table(slide, len(d10a) + 1, 3, Inches(0.5), Inches(1.55), Inches(6.2), Inches(2.8))
for ci, h in enumerate(h10a):
    t10a.columns[ci].width = col10a[ci]
    set_cell_text(t10a.cell(0, ci), h, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    set_cell_fill(t10a.cell(0, ci), ACCENT_BLUE)
for ri, row in enumerate(d10a):
    for ci, val in enumerate(row):
        set_cell_text(t10a.cell(ri + 1, ci), val, font_size=9, alignment=PP_ALIGN.CENTER)
        if ri % 2 == 1:
            set_cell_fill(t10a.cell(ri + 1, ci), ROW_ALT)

add_text_box(slide, Inches(7.0), Inches(1.1), Inches(5.5), Inches(0.5),
             "KV — 关键流程", font_size=18, color=ACCENT_ORANGE, bold=True)

h10b = ["关键流程", "灰白盒关注点", "测试方法"]
col10b = [Inches(1.8), Inches(2.3), Inches(1.9)]
d10b = [
    ["KV Put / Get", "边界键值大小\n并发读写竞态", "边界值 + 协议分析仪"],
    ["KV Delete", "并发访问\n空间回收", "并发路径 + 覆盖率"],
    ["KV List / Iterate", "迭代中数据变更\n游标有效性", "数据流 + 异常注入"],
    ["异常场景", "空间满/闪断\n多控切换", "SFMEA + 网损仪"],
]
t10b = add_table(slide, len(d10b) + 1, 3, Inches(7.0), Inches(1.55), Inches(6.0), Inches(2.8))
for ci, h in enumerate(h10b):
    t10b.columns[ci].width = col10b[ci]
    set_cell_text(t10b.cell(0, ci), h, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    set_cell_fill(t10b.cell(0, ci), ACCENT_ORANGE)
for ri, row in enumerate(d10b):
    for ci, val in enumerate(row):
        set_cell_text(t10b.cell(ri + 1, ci), val, font_size=9, alignment=PP_ALIGN.CENTER)
        if ri % 2 == 1:
            set_cell_fill(t10b.cell(ri + 1, ci), ROW_ALT)

example_box = add_rounded_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.6),
                                fill_color=RGBColor(0xEF, 0xF6, 0xFF), line_color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(0.4),
             "灰盒用例示例：NVMe-oF TCP 连接建立 + 网络闪断 多函数交汇场景",
             font_size=14, color=ACCENT_BLUE, bold=True)

example_content = (
    "关联函数：nvme_tcp_connect() + handle_network_flap() + cleanup_resources()\n"
    "前置条件：NVMe-oF TCP 连接已配置，Target端就绪\n"
    "测试步骤：① 发起连接建立 → ② 在握手过程中注入网络闪断 → ③ 观察连接状态和资源释放\n"
    "预期失败（可接受）：连接失败、返回超时错误码、自动重试\n"
    "不可接受结果：进程崩溃、控制器下电、资源泄漏、不可恢复状态"
)
add_text_box(slide, Inches(0.8), Inches(5.15), Inches(11.5), Inches(1.9),
             example_content, font_size=11, color=MID_TEXT)

# ============================================================
# SLIDE 11 — Milestones
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "计划", 11)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "项目里程碑计划（6个月）", font_size=28, color=DARK_TEXT, bold=True)

milestones = [
    ("M0\n启动", "4月\n1-2周", ACCENT_BLUE,
     ["团队组建", "基线数据采集", "特性选定", "培训计划"]),
    ("M1\n流程", "4月3周\n-5月2周", ACCENT_BLUE,
     ["流程文档", "出入口标准", "工具就绪", "模板输出"]),
    ("M2\n试点", "5月3周\n-6月", ACCENT_TEAL,
     ["NVMe-oF分析", "KV分析启动", "灰盒用例", "AI试跑"]),
    ("M3\n执行", "7月", ACCENT_TEAL,
     ["DT用例执行", "覆盖率跟踪", "问题闭环", "中期评审"]),
    ("M4\n扩展", "8月", ACCENT_ORANGE,
     ["KV推广", "AI深化", "GrayScope验证", "知识库"]),
    ("M5\n总结", "9月", ACCENT_ORANGE,
     ["度量汇总", "经验沉淀", "改进报告", "下阶段规划"]),
]

timeline_y = Inches(1.4)
bar_y = timeline_y + Inches(1.6)
add_shape(slide, Inches(0.8), bar_y, Inches(11.7), Inches(0.08), fill_color=ACCENT_BLUE)

for i, (name, period, color, items) in enumerate(milestones):
    x = Inches(0.8 + i * 2.0)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), bar_y - Inches(0.12),
                                  Inches(0.32), Inches(0.32))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    add_text_box(slide, x, timeline_y, Inches(1.9), Inches(0.6),
                 name, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, timeline_y + Inches(0.7), Inches(1.9), Inches(0.5),
                 period, font_size=10, color=MID_TEXT, alignment=PP_ALIGN.CENTER)

    card = add_rounded_rect(slide, x, bar_y + Inches(0.5), Inches(1.9), Inches(2.0),
                            fill_color=LIGHT_BG, line_color=color)
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.15), bar_y + Inches(0.6 + j * 0.4),
                     Inches(1.6), Inches(0.35),
                     f"▸ {item}", font_size=10, color=MID_TEXT)

key_box = add_rounded_rect(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.2),
                            fill_color=RGBColor(0xEF, 0xF6, 0xFF), line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(6.1), Inches(11), Inches(0.3),
             "关键检查点", font_size=13, color=ACCENT_BLUE, bold=True)
checkpoints = [
    "M1 检查点：流程文档评审通过，工具环境就绪",
    "M3 检查点：中期评审 — 数据有正向趋势，调整后续策略",
    "M5 检查点：最终评审 — 30% 改进目标达成情况，下阶段决策",
]
for j, cp in enumerate(checkpoints):
    add_text_box(slide, Inches(1.3), Inches(6.4 + j * 0.25), Inches(10.5), Inches(0.25),
                 cp, font_size=10, color=MID_TEXT)

# ============================================================
# SLIDE 12 — Team Organization
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "组织", 12)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "团队组织与分工", font_size=28, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
             "3-5 人核心团队，从各组抽调", font_size=14, color=MID_TEXT)

pm_box = add_rounded_rect(slide, Inches(5.0), Inches(1.4), Inches(3.3), Inches(1.0),
                           fill_color=ACCENT_BLUE)
add_text_box(slide, Inches(5.2), Inches(1.5), Inches(2.9), Inches(0.35),
             "项目负责人", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.2), Inches(1.9), Inches(2.9), Inches(0.35),
             "整体推进 | 里程碑管控 | 汇报", font_size=10, color=RGBColor(0xBF, 0xDB, 0xFE),
             alignment=PP_ALIGN.CENTER)

roles = [
    ("协议测试专家", "1-2人 | NVMe/KV测试组",
     ["协议抓包分析", "协议规范梳理", "异常场景设计", "DT用例评审"], ACCENT_TEAL),
    ("代码分析专家", "1人 | 白盒/精准测试组",
     ["代码流程走读", "调用链分析", "覆盖率分析", "灰盒用例设计"], ACCENT_ORANGE),
    ("工具 & AI 专家", "1人 | 工具/自动化组",
     ["精准测试工具", "AI模型对接", "GrayScope验证", "工具链搭建"], RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (title, subtitle, duties, color) in enumerate(roles):
    x = Inches(0.8 + i * 4.2)
    y = Inches(3.0)

    line_shape = add_shape(slide, Inches(6.65), Inches(2.4), Inches(0.03), Inches(0.6), fill_color=BORDER_GRAY)

    card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(2.8),
                            fill_color=LIGHT_BG, line_color=color)
    add_shape(slide, x, y, Inches(3.8), Inches(0.06), fill_color=color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(0.3),
                 subtitle, font_size=10, color=MID_TEXT)
    for j, d in enumerate(duties):
        add_text_box(slide, x + Inches(0.2), y + Inches(1.0 + j * 0.35),
                     Inches(3.4), Inches(0.3),
                     f"▸ {d}", font_size=11, color=DARK_TEXT)

collab_box = add_rounded_rect(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(1.1),
                               fill_color=LIGHT_BG, line_color=BORDER_GRAY)
add_text_box(slide, Inches(1.1), Inches(6.2), Inches(11), Inches(0.3),
             "协作机制", font_size=13, color=DARK_TEXT, bold=True)
collabs = [
    "每周 1 次站会（30分钟）— 对齐进度和问题",
    "每个里程碑 1 次评审会 — 检查交付物和数据",
    "人员最低投入占比 ≥ 30%，关键里程碑设 Hard Deadline",
]
for j, c in enumerate(collabs):
    add_text_box(slide, Inches(1.3), Inches(6.5 + j * 0.2), Inches(10.5), Inches(0.2),
                 c, font_size=10, color=MID_TEXT)

# ============================================================
# SLIDE 13 — Metrics Dashboard
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "度量", 13)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "度量看板设计", font_size=28, color=DARK_TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
             "结果指标 (R) / 过程指标 (P) / 能力指标 (C) 三维模型", font_size=14, color=MID_TEXT)

dims = [
    ("R 结果指标", "用数据证明改进效果",
     [("R1", "异常分支类漏测下降率", "≥30%"),
      ("R2", "并发类漏测下降率", "≥30%"),
      ("R3", "协议报文类漏测下降率", "≥30%")],
     ACCENT_RED, RGBColor(0xFE, 0xF2, 0xF2)),
    ("P 过程指标", "确保活动被执行",
     [("P1", "灰白盒活动执行率", "≥90%"),
      ("P2", "覆盖率增量", "可量化"),
      ("P3", "灰盒用例产出数", "≥30条")],
     ACCENT_TEAL, RGBColor(0xEC, 0xFD, 0xF5)),
    ("C 能力指标", "验证工具和AI的价值",
     [("C1", "AI辅助分析准确率", "≥70%"),
      ("C2", "GrayScope验证达标率", "≥70%"),
      ("C3", "调用链→场景转化", "≥20条")],
     ACCENT_ORANGE, RGBColor(0xFF, 0xF7, 0xED)),
]

for i, (title, subtitle, metrics, color, bg) in enumerate(dims):
    x = Inches(0.6 + i * 4.2)
    card = add_rounded_rect(slide, x, Inches(1.4), Inches(3.9), Inches(4.5),
                            fill_color=bg, line_color=color)
    add_shape(slide, x, Inches(1.4), Inches(3.9), Inches(0.06), fill_color=color)
    add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(3.5), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(3.5), Inches(0.3),
                 subtitle, font_size=11, color=MID_TEXT)

    for j, (mid, mname, mtarget) in enumerate(metrics):
        my = Inches(2.5 + j * 1.1)
        m_card = add_rounded_rect(slide, x + Inches(0.15), my, Inches(3.6), Inches(0.9),
                                   fill_color=WHITE, line_color=BORDER_GRAY)
        add_text_box(slide, x + Inches(0.3), my + Inches(0.05), Inches(0.6), Inches(0.3),
                     mid, font_size=10, color=color, bold=True)
        add_text_box(slide, x + Inches(0.9), my + Inches(0.05), Inches(2.6), Inches(0.3),
                     mname, font_size=11, color=DARK_TEXT)
        target_box = add_rounded_rect(slide, x + Inches(2.2), my + Inches(0.45),
                                       Inches(1.4), Inches(0.35), fill_color=color)
        add_text_box(slide, x + Inches(2.2), my + Inches(0.47), Inches(1.4), Inches(0.3),
                     mtarget, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tracking = add_rounded_rect(slide, Inches(0.6), Inches(6.15), Inches(12.1), Inches(1.1),
                             fill_color=LIGHT_BG, line_color=BORDER_GRAY)
add_text_box(slide, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.3),
             "度量采集节奏", font_size=13, color=DARK_TEXT, bold=True)
rhythm = (
    "▸ R 指标：每版本从缺陷系统按标签分类统计，与基线对比\n"
    "▸ P 指标：每版本从测试报告和用例系统采集\n"
    "▸ C 指标：M4-M5 阶段通过人工评审采集"
)
add_text_box(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.6),
             rhythm, font_size=10, color=MID_TEXT)

# ============================================================
# SLIDE 14 — Risks
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "风险", 14)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "风险识别与应对", font_size=28, color=DARK_TEXT, bold=True)

risks = [
    ("团队代码分析能力不足", "高", "中",
     "M0针对性培训；选有代码基础成员；\n前期代码分析专家带教"),
    ("AI模型效果不达预期", "中", "中",
     "AI定位为\"辅助\"非\"依赖\"；\n方案一二独立可运作；迭代优化prompt"),
    ("GrayScope成熟度不够", "中", "高",
     "明确为\"验证性\"引入；\n不纳入不影响核心目标"),
    ("试点特性版本计划变动", "中", "中",
     "紧跟版本计划调整里程碑；\n保持灵活分析窗口"),
    ("基线数据统计口径不一", "高", "低",
     "M0明确统计口径和分类标准；\n与QA对齐缺陷标签规范"),
    ("人员投入不足（兼职）", "中", "中",
     "明确最低投入≥30%；\n关键里程碑设Hard Deadline"),
]

risk_colors = {"高": ACCENT_RED, "中": ACCENT_ORANGE, "低": ACCENT_GREEN}

headers14 = ["风险描述", "影响", "概率", "应对措施"]
col14 = [Inches(3.0), Inches(0.8), Inches(0.8), Inches(6.8)]
table14 = add_table(slide, len(risks) + 1, 4,
                    Inches(0.5), Inches(1.2), Inches(11.4), Inches(5.6))
for ci, h in enumerate(headers14):
    table14.columns[ci].width = col14[ci]
    set_cell_text(table14.cell(0, ci), h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    set_cell_fill(table14.cell(0, ci), ACCENT_RED)

for ri, (desc, impact, prob, measure) in enumerate(risks):
    set_cell_text(table14.cell(ri + 1, 0), desc, font_size=10, bold=True)
    set_cell_text(table14.cell(ri + 1, 1), impact, font_size=10, alignment=PP_ALIGN.CENTER,
                  color=risk_colors.get(impact, MID_TEXT), bold=True)
    set_cell_text(table14.cell(ri + 1, 2), prob, font_size=10, alignment=PP_ALIGN.CENTER,
                  color=risk_colors.get(prob, MID_TEXT), bold=True)
    set_cell_text(table14.cell(ri + 1, 3), measure, font_size=10)
    if ri % 2 == 1:
        for ci in range(4):
            set_cell_fill(table14.cell(ri + 1, ci), ROW_ALT)

# ============================================================
# SLIDE 15 — Expected Benefits
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "收益", 15)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "预期收益与展望", font_size=28, color=DARK_TEXT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
             "短期收益（项目周期内）", font_size=18, color=ACCENT_BLUE, bold=True)

benefits = [
    ("漏测下降 ≥30%", "异常分支/并发/协议三类", ACCENT_RED),
    ("专业能力提升", "团队具备灰白盒设计能力", ACCENT_TEAL),
    ("流程标准建立", "可复用的灰白盒流程文档", ACCENT_BLUE),
    ("用例资产沉淀", "≥30 条高质量灰盒用例", ACCENT_ORANGE),
    ("AI 可行性验证", "精准测试+AI 验证报告", RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (title, desc, color) in enumerate(benefits):
    x = Inches(0.8 + (i % 3) * 4.0)
    y = Inches(1.65 + (i // 3) * 1.3)
    card = add_rounded_rect(slide, x, y, Inches(3.7), Inches(1.1),
                            fill_color=LIGHT_BG, line_color=color)
    add_shape(slide, x, y, Inches(0.1), Inches(1.1), fill_color=color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(3.2), Inches(0.4),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.55), Inches(3.2), Inches(0.4),
                 desc, font_size=11, color=MID_TEXT)

add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.5),
             "中长期推广路线", font_size=18, color=ACCENT_TEAL, bold=True)

roadmap = [
    ("当前", "试点 2 个特性\nNVMe-oF TCP + KV"),
    ("推广", "覆盖更多协议\niSCSI / FC / NAS..."),
    ("标准化", "纳入盘古版本\n交付标准流程"),
    ("智能化", "AI 工具链\n成为标配"),
]

for i, (stage, desc) in enumerate(roadmap):
    x = Inches(0.8 + i * 3.1)
    y = Inches(4.9)
    box = add_rounded_rect(slide, x, y, Inches(2.6), Inches(1.3),
                            fill_color=LIGHT_BG, line_color=ACCENT_TEAL)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.4), Inches(0.35),
                 stage, font_size=14, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.5), Inches(2.4), Inches(0.7),
                 desc, font_size=11, color=MID_TEXT, alignment=PP_ALIGN.CENTER)
    if i < 3:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.65), y + Inches(0.4),
                                      Inches(0.4), Inches(0.35))
        arr.fill.solid()
        arr.fill.fore_color.rgb = ACCENT_TEAL
        arr.line.fill.background()

value_box = add_rounded_rect(slide, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.85),
                              fill_color=RGBColor(0xEF, 0xF6, 0xFF), line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(6.5), Inches(11), Inches(0.3),
             "对外展示价值", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.1), Inches(6.8), Inches(11), Inches(0.4),
             "测试专业能力（技术深度） | 方法论输出（可推广） | 工具链创新（前沿探索） | 质量实效（30%硬指标）",
             font_size=11, color=MID_TEXT)

# ============================================================
# SLIDE 16 — Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "行动", 16)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "下一步行动", font_size=28, color=DARK_TEXT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
             "立即行动（本周/下周）", font_size=18, color=ACCENT_BLUE, bold=True)

headers16 = ["#", "行动项", "责任人", "完成时间"]
col16 = [Inches(0.5), Inches(5.5), Inches(2.5), Inches(2.0)]
data16 = [
    ["1", "确认项目团队成员名单", "项目负责人", "本周"],
    ["2", "从缺陷系统导出基线数据", "项目负责人", "本周"],
    ["3", "确认 NVMe-oF TCP / KV 关键流程清单", "协议测试专家", "下周"],
    ["4", "搭建精准测试工具环境", "工具 & AI 专家", "下周"],
    ["5", "制定团队培训计划", "项目负责人", "下周"],
]

table16 = add_table(slide, len(data16) + 1, 4,
                    Inches(0.5), Inches(1.6), Inches(10.5), Inches(2.8))
for ci, h in enumerate(headers16):
    table16.columns[ci].width = col16[ci]
    set_cell_text(table16.cell(0, ci), h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    set_cell_fill(table16.cell(0, ci), ACCENT_BLUE)
for ri, row in enumerate(data16):
    for ci, val in enumerate(row):
        set_cell_text(table16.cell(ri + 1, ci), val, font_size=11, alignment=PP_ALIGN.CENTER)
        if ri % 2 == 1:
            for c in range(4):
                set_cell_fill(table16.cell(ri + 1, c), ROW_ALT)

add_text_box(slide, Inches(0.8), Inches(4.7), Inches(5), Inches(0.5),
             "需要决策 / 支持", font_size=18, color=ACCENT_ORANGE, bold=True)

decisions = [
    ("各组抽调人员的投入比例确认", "建议 ≥ 30%"),
    ("NVMe-oF TCP / KV 关键流程选择", "建议各选 3-5 个"),
    ("公司 AI 模型接入权限和技术支持", "联系AI平台团队"),
]

for i, (item, note) in enumerate(decisions):
    y = Inches(5.3 + i * 0.7)
    card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.55),
                            fill_color=RGBColor(0xFF, 0xF7, 0xED), line_color=ACCENT_ORANGE)
    add_text_box(slide, Inches(1.1), y + Inches(0.08), Inches(7), Inches(0.35),
                 f"▸ {item}", font_size=12, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(8.5), y + Inches(0.08), Inches(3.5), Inches(0.35),
                 note, font_size=11, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.RIGHT)

# ============================================================
# Save
# ============================================================
output_path = "/Volumes/Media/Coding/docs/plans/盘古平台灰白盒测试系统建设项目计划.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
